# Document Processing Service for AI Study Assistant
# Handles file uploads, text extraction, and content analysis

from typing import List, Dict, Any, Optional
import asyncio
import hashlib
import json
import logging
from datetime import datetime
from pathlib import Path
import aiofiles
from fastapi import UploadFile
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    import tiktoken
except ImportError:
    tiktoken = None

from services.gemini_service import get_gemini_service

# Configure logging
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """
    Revolutionary document processing system that extracts and analyzes content
    from uploaded files using LLaMA 3.2 for intelligent processing
    """
    
    def __init__(self):
        self.supported_formats = {'.pdf', '.pptx', '.docx', '.txt'}
        self.upload_dir = Path("uploads")
        self.upload_dir.mkdir(exist_ok=True)
        self.tokenizer = tiktoken.encoding_for_model("gpt-3.5-turbo") if tiktoken else None  # For token counting
    
    async def process_document(self, file: UploadFile, user_id: str) -> Dict[str, Any]:
        """
        Alias for process_upload to maintain compatibility
        """
        return await self.process_upload(file, user_id)
    
    async def process_document_streaming(self, file: UploadFile, user_id: str, progress_callback=None):
        """
        Process document with page-by-page streaming results
        Shows analysis for each page as it's processed instead of waiting for entire document
        """
        try:
            # Validate file
            if not self._is_supported_format(file.filename):
                raise ValueError(f"Unsupported file format. Supported: {self.supported_formats}")
            
            # Save file
            file_path = await self._save_file(file, user_id)
            
            if progress_callback:
                await progress_callback("📁 File saved, starting content extraction...")
            
            # Extract content with page-by-page processing
            content = await self._extract_content_streaming(file_path, file.filename, progress_callback)
            
            if progress_callback:
                await progress_callback("🎯 Document processing complete!")
            
            return content
            
        except Exception as e:
            raise Exception(f"Streaming document processing failed: {str(e)}")
    
    async def process_upload(self, file: UploadFile, user_id: str) -> Dict[str, Any]:
        """
        Main processing pipeline for uploaded documents
        """
        try:
            # Validate file
            if not self._is_supported_format(file.filename):
                raise ValueError(f"Unsupported file format. Supported: {self.supported_formats}")
            
            # Save file
            file_path = await self._save_file(file, user_id)
            
            # Extract content
            content = await self._extract_content(file_path, file.filename)
            
            # Analyze with LLaMA 3.2
            analysis = await self._analyze_content_with_llama(content, file.filename)
            
            # Generate study materials
            study_materials = await self._generate_study_materials(content, analysis, file.filename)
            
            # Save to database
            document_data = await self._save_to_database(
                user_id, file.filename, file_path, content, analysis, study_materials
            )
            
            return document_data
            
        except Exception as e:
            raise Exception(f"Document processing failed: {str(e)}")
    
    def _is_supported_format(self, filename: str) -> bool:
        """Check if file format is supported"""
        return Path(filename).suffix.lower() in self.supported_formats
    
    async def _save_file(self, file: UploadFile, user_id: str) -> Path:
        """Save uploaded file to disk"""
        # Create unique filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        file_hash = hashlib.md5(f"{user_id}_{timestamp}_{file.filename}".encode()).hexdigest()[:8]
        safe_filename = f"{timestamp}_{file_hash}_{file.filename}"
        
        file_path = self.upload_dir / user_id / safe_filename
        file_path.parent.mkdir(parents=True, exist_ok=True)
        
        # Save file
        async with aiofiles.open(file_path, 'wb') as f:
            content = await file.read()
            await f.write(content)
        
        return file_path
    
    async def _extract_content(self, file_path: Path, filename: str) -> Dict[str, Any]:
        """Extract text content from different file formats"""
        file_extension = file_path.suffix.lower()
        
        if file_extension == '.pdf':
            return await self._extract_pdf_content(file_path)
        elif file_extension == '.pptx':
            return await self._extract_pptx_content(file_path)
        elif file_extension == '.docx':
            return await self._extract_docx_content(file_path)
        elif file_extension == '.txt':
            return await self._extract_txt_content(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    
    async def _extract_pdf_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from PDF files"""
        content = {
            'pages': [],
            'total_pages': 0,
            'text': '',
            'metadata': {}
        }
        
        try:
            if not PyPDF2:
                raise Exception("PyPDF2 not installed")
                
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                content['total_pages'] = len(pdf_reader.pages)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    content['pages'].append({
                        'page_number': page_num + 1,
                        'text': page_text,
                        'word_count': len(page_text.split())
                    })
                    content['text'] += page_text + "\n\n"
                
                # Extract metadata
                if pdf_reader.metadata:
                    content['metadata'] = {
                        'title': pdf_reader.metadata.get('/Title', ''),
                        'author': pdf_reader.metadata.get('/Author', ''),
                        'subject': pdf_reader.metadata.get('/Subject', ''),
                        'creator': pdf_reader.metadata.get('/Creator', '')
                    }
        
        except Exception as e:
            raise Exception(f"PDF extraction failed: {str(e)}")
        
        return content
    
    async def _extract_pptx_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from PowerPoint files"""
        content = {
            'slides': [],
            'total_slides': 0,
            'text': '',
            'metadata': {}
        }
        
        try:
            if not Presentation:
                raise Exception("python-pptx not installed")
            prs = Presentation(file_path)
            content['total_slides'] = len(prs.slides)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_content = {
                    'slide_number': slide_num + 1,
                    'title': '',
                    'text': '',
                    'bullet_points': []
                }
                
                # Extract text from shapes
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            if not slide_content['title'] and len(shape.text) < 100:
                                slide_content['title'] = shape.text.strip()
                            else:
                                slide_text.append(shape.text.strip())
                                # Check if it's a bullet point
                                if shape.text.strip().startswith(('•', '-', '*')):
                                    slide_content['bullet_points'].append(shape.text.strip())
                
                slide_content['text'] = '\n'.join(slide_text)
                content['slides'].append(slide_content)
                content['text'] += f"Slide {slide_num + 1}: {slide_content['title']}\n{slide_content['text']}\n\n"
            
            # Extract presentation metadata
            if hasattr(prs.core_properties, 'title'):
                content['metadata'] = {
                    'title': prs.core_properties.title or '',
                    'author': prs.core_properties.author or '',
                    'subject': prs.core_properties.subject or '',
                    'created': str(prs.core_properties.created) if prs.core_properties.created else ''
                }
        
        except Exception as e:
            raise Exception(f"PowerPoint extraction failed: {str(e)}")
        
        return content
    
    async def _extract_docx_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from Word documents"""
        content = {
            'paragraphs': [],
            'text': '',
            'metadata': {}
        }
        
        try:
            if not Document:
                raise Exception("python-docx not installed")
            doc = Document(file_path)
            
            for para_num, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    content['paragraphs'].append({
                        'paragraph_number': para_num + 1,
                        'text': paragraph.text.strip(),
                        'style': paragraph.style.name if paragraph.style else 'Normal'
                    })
                    content['text'] += paragraph.text.strip() + '\n\n'
            
            # Extract document properties
            content['metadata'] = {
                'title': doc.core_properties.title or '',
                'author': doc.core_properties.author or '',
                'subject': doc.core_properties.subject or '',
                'created': str(doc.core_properties.created) if doc.core_properties.created else ''
            }
        
        except Exception as e:
            raise Exception(f"Word document extraction failed: {str(e)}")
        
        return content
    
    async def _extract_txt_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from text files"""
        content = {
            'text': '',
            'lines': [],
            'metadata': {}
        }
        
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as file:
                text = await file.read()
                content['text'] = text
                content['lines'] = text.split('\n')
                content['metadata'] = {
                    'line_count': len(content['lines']),
                    'character_count': len(text),
                    'word_count': len(text.split())
                }
        
        except Exception as e:
            raise Exception(f"Text file extraction failed: {str(e)}")
        
        return content
    
    async def _extract_content_streaming(self, file_path: Path, filename: str, progress_callback=None) -> Dict[str, Any]:
        """
        Extract content with page-by-page Gemini analysis and real-time updates
        """
        file_extension = file_path.suffix.lower()
        
        if file_extension == '.pdf':
            return await self._extract_pdf_streaming(file_path, filename, progress_callback)
        elif file_extension == '.pptx':
            return await self._extract_pptx_streaming(file_path, filename, progress_callback)
        elif file_extension == '.docx':
            return await self._extract_docx_streaming(file_path, filename, progress_callback)
        elif file_extension == '.txt':
            return await self._extract_text_streaming(file_path, filename, progress_callback)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    
    async def _extract_pdf_streaming(self, file_path: Path, filename: str, progress_callback=None) -> Dict[str, Any]:
        """Extract PDF content page-by-page with immediate Gemini analysis"""
        if not PyPDF2:
            raise Exception("PyPDF2 not installed")
        
        content = {
            'type': 'pdf',
            'filename': filename,
            'text': '',
            'pages': [],
            'analysis_results': [],
            'flashcards': [],
            'total_pages': 0,
            'processed_pages': 0
        }
        
        gemini_service = get_gemini_service()
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            content['total_pages'] = len(pdf_reader.pages)
            
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                
                # Add page info
                page_info = {
                    'page_number': page_num + 1,
                    'text': page_text,
                    'word_count': len(page_text.split())
                }
                content['pages'].append(page_info)
                content['text'] += page_text + "\n\n"
                content['processed_pages'] = page_num + 1
                
                if progress_callback:
                    await progress_callback(f"📄 Extracting page {page_num + 1}/{len(pdf_reader.pages)}...")
                
                # Analyze page immediately with Gemini if it has enough content
                if len(page_text.strip()) > 50:  # Only analyze pages with substantial content
                    try:
                        if progress_callback:
                            await progress_callback(f"🤖 Gemini analyzing page {page_num + 1}...")
                        
                        # Get page analysis
                        page_analysis = await gemini_service.analyze_content(page_text, f"{filename} - Page {page_num + 1}")
                        page_analysis['page_number'] = page_num + 1
                        content['analysis_results'].append(page_analysis)
                        
                        # Generate flashcards for this page
                        page_flashcards = await gemini_service.generate_flashcards(page_text, f"{filename} - Page {page_num + 1}", 3)
                        for card in page_flashcards:
                            card['page_number'] = page_num + 1
                        content['flashcards'].extend(page_flashcards)
                        
                        if progress_callback:
                            await progress_callback(f"✅ Page {page_num + 1} analyzed - {len(page_flashcards)} flashcards created")
                        
                    except Exception as e:
                        logger.error(f"Error analyzing page {page_num + 1}: {e}")
                        if progress_callback:
                            await progress_callback(f"⚠️ Page {page_num + 1} analysis skipped")
        
        # Generate final summary from all content
        if progress_callback:
            await progress_callback("📝 Creating final document summary...")
        
        try:
            final_summary = await gemini_service.generate_summary(content['text'], filename)
            content['final_summary'] = final_summary
        except Exception as e:
            logger.error(f"Error generating final summary: {e}")
            content['final_summary'] = "Summary generation failed"
        
        return content
    
    async def _extract_pptx_streaming(self, file_path: Path, filename: str, progress_callback=None) -> Dict[str, Any]:
        """Extract PowerPoint content slide-by-slide with immediate Gemini analysis"""
        if not Presentation:
            raise Exception("python-pptx not installed")
        
        prs = Presentation(file_path)
        content = {
            'type': 'presentation',
            'filename': filename,
            'text': '',
            'slides': [],
            'analysis_results': [],
            'flashcards': [],
            'total_slides': len(prs.slides),
            'processed_slides': 0
        }
        
        gemini_service = get_gemini_service()
        
        for slide_num, slide in enumerate(prs.slides):
            slide_content = {
                'slide_number': slide_num + 1,
                'title': '',
                'text': '',
                'bullet_points': []
            }
            
            # Extract text from shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        if not slide_content['title'] and len(shape.text) < 100:
                            slide_content['title'] = shape.text.strip()
                        else:
                            slide_text.append(shape.text.strip())
                            # Check if it's a bullet point
                            if shape.text.strip().startswith(('•', '-', '*')):
                                slide_content['bullet_points'].append(shape.text.strip())
            
            slide_content['text'] = '\n'.join(slide_text)
            content['slides'].append(slide_content)
            content['text'] += f"Slide {slide_num + 1}: {slide_content['title']}\n{slide_content['text']}\n\n"
            content['processed_slides'] = slide_num + 1
            
            if progress_callback:
                await progress_callback(f"📊 Processing slide {slide_num + 1}/{len(prs.slides)}: {slide_content['title'][:50]}...")
            
            # Analyze slide immediately with Gemini
            slide_full_text = f"{slide_content['title']}\n{slide_content['text']}"
            if len(slide_full_text.strip()) > 20:
                try:
                    if progress_callback:
                        await progress_callback(f"🤖 Gemini analyzing slide {slide_num + 1}...")
                    
                    # Get slide analysis
                    slide_analysis = await gemini_service.analyze_content(slide_full_text, f"{filename} - Slide {slide_num + 1}")
                    slide_analysis['slide_number'] = slide_num + 1
                    slide_analysis['slide_title'] = slide_content['title']
                    content['analysis_results'].append(slide_analysis)
                    
                    # Generate flashcards for this slide
                    slide_flashcards = await gemini_service.generate_flashcards(slide_full_text, f"{filename} - Slide {slide_num + 1}", 2)
                    for card in slide_flashcards:
                        card['slide_number'] = slide_num + 1
                        card['slide_title'] = slide_content['title']
                    content['flashcards'].extend(slide_flashcards)
                    
                    if progress_callback:
                        await progress_callback(f"✅ Slide {slide_num + 1} analyzed - {len(slide_flashcards)} flashcards created")
                    
                except Exception as e:
                    logger.error(f"Error analyzing slide {slide_num + 1}: {e}")
                    if progress_callback:
                        await progress_callback(f"⚠️ Slide {slide_num + 1} analysis skipped")
        
        # Generate final summary
        if progress_callback:
            await progress_callback("📝 Creating final presentation summary...")
        
        try:
            final_summary = await gemini_service.generate_summary(content['text'], filename)
            content['final_summary'] = final_summary
        except Exception as e:
            logger.error(f"Error generating final summary: {e}")
            content['final_summary'] = "Summary generation failed"
        
        return content
    
    async def _extract_text_streaming(self, file_path: Path, filename: str, progress_callback=None) -> Dict[str, Any]:
        """Extract text file with immediate Gemini analysis"""
        content = {
            'type': 'text',
            'filename': filename,
            'text': '',
            'analysis_results': [],
            'flashcards': [],
            'lines': []
        }
        
        if progress_callback:
            await progress_callback("📄 Reading text file...")
        
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as file:
                text = await file.read()
                content['text'] = text
                content['lines'] = text.split('\n')
                
                if progress_callback:
                    await progress_callback("🤖 Gemini analyzing text content...")
                
                # Analyze the text content
                gemini_service = get_gemini_service()
                
                try:
                    # Get content analysis
                    analysis = await gemini_service.analyze_content(text, filename)
                    content['analysis_results'] = [analysis]
                    
                    # Generate flashcards
                    flashcards = await gemini_service.generate_flashcards(text, filename, 5)
                    content['flashcards'] = flashcards
                    
                    # Generate summary
                    summary = await gemini_service.generate_summary(text, filename)
                    content['final_summary'] = summary
                    
                    if progress_callback:
                        await progress_callback(f"✅ Analysis complete - {len(flashcards)} flashcards created")
                        
                except Exception as e:
                    logger.error(f"Error analyzing text: {e}")
                    if progress_callback:
                        await progress_callback("⚠️ Analysis failed")
                
        except Exception as e:
            raise Exception(f"Text file extraction failed: {str(e)}")
        
        return content
    
    async def _analyze_content_with_llama(self, content: Dict[str, Any], filename: str) -> Dict[str, Any]:
        """
        Use LLaMA 3.2 to analyze and understand the content
        This is where the magic happens!
        """
        text = content.get('text', '')
        
        # Debug: Check if text is actually a string
        if not isinstance(text, str):
            logger.error(f"❌ Expected string but got {type(text)}: {text}")
            text = str(text) if text else ''
        
        # Create analysis prompt for LLaMA
        analysis_prompt = f"""
        Analyze this educational content and provide a comprehensive analysis:

        CONTENT TO ANALYZE:
        {text[:4000]}  # Truncate for token limits

        Please provide:
        1. SUBJECT_AREA: What subject/field is this content about?
        2. DIFFICULTY_LEVEL: Rate from 1-10 (beginner to expert)
        3. KEY_CONCEPTS: List the main concepts/topics covered
        4. LEARNING_OBJECTIVES: What should students learn from this?
        5. PREREQUISITE_KNOWLEDGE: What should students know beforehand?
        6. CONTENT_STRUCTURE: How is the information organized?
        7. TEACHING_APPROACH: How would you teach this content?
        8. ASSESSMENT_IDEAS: What questions would test understanding?

        Respond in JSON format.
        """
        
        try:
            # Get analysis from Gemini API (FAST!)
            gemini_service = get_gemini_service()
            analysis = await gemini_service.analyze_content(text, filename)
            return analysis
            
        except Exception as e:
            # Fallback analysis if LLaMA fails
            return {
                'subject_area': 'General',
                'difficulty_level': 5,
                'key_concepts': self._extract_keywords(text),
                'learning_objectives': ['Understand the main concepts presented'],
                'prerequisite_knowledge': ['Basic understanding of the subject'],
                'content_structure': 'Sequential presentation',
                'teaching_approach': 'Step-by-step explanation',
                'assessment_ideas': ['Multiple choice questions', 'Short answer questions'],
                'error': f"LLaMA analysis failed: {str(e)}",
                'processing_metadata': {
                    'filename': filename,
                    'processed_at': datetime.utcnow().isoformat(),
                    'content_length': len(text),
                    'fallback_analysis': True
                }
            }
    
    def _extract_keywords(self, text: str) -> List[str]:
        """Simple keyword extraction fallback"""
        words = text.lower().split()
        # Filter out common words and get unique keywords
        common_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'is', 'are', 'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could', 'should', 'may', 'might', 'must', 'can', 'this', 'that', 'these', 'those'}
        keywords = [word for word in set(words) if len(word) > 3 and word not in common_words]
        return keywords[:20]  # Top 20 keywords
    
    async def _generate_study_materials(self, content: Dict[str, Any], analysis: Dict[str, Any], filename: str) -> Dict[str, Any]:
        """Generate study materials using LLaMA 3.2"""
        text = content.get('text', '')
        
        # Generate summary
        summary_prompt = f"""
        Create a comprehensive summary of this educational content:
        
        {text[:3000]}
        
        Provide a clear, concise summary that captures the main points and key information.
        Make it suitable for study and review purposes.
        """
        
        gemini_service = get_gemini_service()
        summary = await gemini_service.generate_summary(text, filename)
        
        # Generate initial flashcards
        flashcard_prompt = f"""
        Create 10 flashcards from this content:
        
        {text[:3000]}
        
        Format each flashcard as:
        Q: [Question]
        A: [Answer]
        
        Focus on key concepts and important facts.
        """
        
        gemini_service = get_gemini_service()
        flashcards_raw = await gemini_service.generate_flashcards(text, filename, 10)
        
        # Convert Gemini flashcards to expected format
        flashcards = self._convert_flashcards_format(flashcards_raw)
        
        return {
            'summary': summary,
            'flashcards': flashcards,
            'generated_at': datetime.utcnow().isoformat(),
            'flashcard_count': len(flashcards)
        }
    
    def _convert_flashcards_format(self, flashcards_list: List[Dict[str, str]]) -> List[Dict[str, str]]:
        """Convert Gemini flashcards format to expected format"""
        converted = []
        for card in flashcards_list:
            converted.append({
                'front': card.get('question', ''),
                'back': card.get('answer', ''),
                'difficulty': 'medium',
                'tags': [],
                'created_at': datetime.now().isoformat()
            })
        return converted
    
    def _parse_flashcards(self, flashcards_text: str) -> List[Dict[str, str]]:
        """Parse flashcards from LLaMA response (legacy method)"""
        flashcards = []
        lines = flashcards_text.split('\n')
        current_question = ""
        current_answer = ""
        
        for line in lines:
            line = line.strip()
            if line.startswith('Q:'):
                if current_question and current_answer:
                    flashcards.append({
                        'front': current_question,
                        'back': current_answer,
                        'difficulty': 'medium',
                        'created_at': datetime.utcnow().isoformat()
                    })
                current_question = line[2:].strip()
                current_answer = ""
            elif line.startswith('A:'):
                current_answer = line[2:].strip()
        
        # Add the last flashcard
        if current_question and current_answer:
            flashcards.append({
                'front': current_question,
                'back': current_answer,
                'difficulty': 'medium',
                'created_at': datetime.utcnow().isoformat()
            })
        
        return flashcards
    
    async def _save_to_database(self, user_id: str, filename: str, file_path: Path, 
                               content: Dict[str, Any], analysis: Dict[str, Any], 
                               study_materials: Dict[str, Any]) -> Dict[str, Any]:
        """Save processed document to in-memory storage"""
        # Simple return without database save
        document_data = {
            'user_id': user_id,
            'filename': filename,
            'file_path': str(file_path),
            'upload_timestamp': datetime.utcnow(),
            'processing_status': 'completed',
            'content': content,
            'analysis': analysis,
            'study_materials': study_materials,
            'metadata': {
                'file_size': file_path.stat().st_size if file_path.exists() else 0,
                'content_type': self._get_content_type(filename),
                'processing_time': datetime.utcnow().isoformat()
            }
        }
        
        # Return without database save
        document_data['_id'] = str(hash(filename + str(datetime.utcnow())))
        
        return document_data
        
        return document_data
    
    def _get_content_type(self, filename: str) -> str:
        """Get content type from filename"""
        extension = Path(filename).suffix.lower()
        types = {
            '.pdf': 'application/pdf',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.txt': 'text/plain'
        }
        return types.get(extension, 'application/octet-stream')

# Global instance
document_processor = DocumentProcessor()